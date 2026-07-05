"""Generate Dokaz_Sales_Deck.pptx — a 7-slide deck to send to a business
prospect. Ends with a sign-up CTA. Same palette as the investor deck
(steel + baby-royal-blue + pink + sea-teal), tuned for the buyer, not the VC.

Voice: name the gap, close it with proof. Every slide leads with something
concrete (the receipt, the shipped feature, the exact price) — nothing
vague, no vision talk.

Run: `python docs/presentations/build_sales_deck.py` — writes Dokaz_Sales_Deck.pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# =============================================================================
# Brand palette — mirrors branding/design-system/tokens/colors.css.
# =============================================================================
BRAND_900 = RGBColor(0x24, 0x3A, 0x78)
BRAND_700 = RGBColor(0x3F, 0x63, 0xC4)
BRAND_500 = RGBColor(0x5B, 0x8D, 0xEF)
BRAND_400 = RGBColor(0x6E, 0x9B, 0xF0)
BRAND_300 = RGBColor(0x8C, 0xB1, 0xF7)

PINK_400  = RGBColor(0xF4, 0x8F, 0xB1)
PINK_300  = RGBColor(0xF9, 0xA8, 0xD4)
TEAL_500  = RGBColor(0x2B, 0xA8, 0x88)
TEAL_400  = RGBColor(0x56, 0xC5, 0x96)
TEAL_300  = RGBColor(0x5F, 0xCB, 0xAC)
DANGER    = RGBColor(0xFF, 0x6B, 0x81)

STEEL_950 = RGBColor(0x11, 0x15, 0x1A)
STEEL_900 = RGBColor(0x1B, 0x20, 0x27)
STEEL_800 = RGBColor(0x25, 0x2C, 0x35)
STEEL_700 = RGBColor(0x33, 0x40, 0x5A)
STEEL_500 = RGBColor(0x5D, 0x6C, 0x80)
STEEL_400 = RGBColor(0x7A, 0x8A, 0x9C)
STEEL_200 = RGBColor(0xC4, 0xCF, 0xDB)

WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0xED, 0xF1, 0xF6)
BODY      = RGBColor(0xD7, 0xDE, 0xE7)
MUTED     = RGBColor(0x9A, 0xA7, 0xB6)

CARD_BG   = RGBColor(0x2E, 0x36, 0x40)

# =============================================================================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

TOTAL = 7


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = STEEL_900
    bg.shadow.inherit = False
    return s


def textbox(slide, left, top, width, height, text, *,
            size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def rounded(slide, left, top, width, height, fill, line=None, line_w=Pt(0.75)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.10
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_w
    shape.shadow.inherit = False
    return shape


def pill(slide, left, top, width, height, text, *, bg=BRAND_500, fg=WHITE, size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    shape.shadow.inherit = False
    tb = shape.text_frame
    tb.margin_left = tb.margin_right = Inches(0.05)
    tb.margin_top = tb.margin_bottom = 0
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    return shape


def page_chrome(slide, eyebrow, title, slide_num, total):
    textbox(slide, Inches(0.6), Inches(0.35), Inches(2), Inches(0.35),
            'Dokaz', size=13, bold=True, color=INK)
    textbox(slide, Inches(12.4), Inches(0.35), Inches(0.6), Inches(0.35),
            f'{slide_num:02d} / {total:02d}',
            size=9, color=STEEL_500, align=PP_ALIGN.RIGHT, font='Consolas')
    textbox(slide, Inches(0.6), Inches(0.85), Inches(6), Inches(0.35),
            eyebrow.upper(), size=10, bold=True, color=PINK_300)
    textbox(slide, Inches(0.6), Inches(1.15), Inches(12), Inches(0.9),
            title, size=30, bold=True, color=INK)
    rect(slide, Inches(0.6), Inches(2.05), Inches(0.6), Inches(0.04), PINK_400)


# =============================================================================
# Slide 1 — Cover
# =============================================================================
s = add_slide()

glow_l = rect(s, Inches(-2), Inches(-1), Inches(6), Inches(5), BRAND_900)
glow_l.line.fill.background()

textbox(s, Inches(0.9), Inches(2.5), Inches(6), Inches(0.9),
        'Dokaz', size=68, bold=True, color=WHITE)
textbox(s, Inches(0.9), Inches(3.5), Inches(9), Inches(0.5),
        'Prove your backups actually restore.',
        size=22, color=PINK_300)
textbox(s, Inches(0.9), Inches(4.1), Inches(11), Inches(0.9),
        'A drill is a real restore in an isolated sandbox on a schedule you set. '
        'Every run produces a signed Proof-of-Recovery PDF your auditor and insurer '
        'accept — with an independent verifier so no one has to trust us.',
        size=13, color=BODY)

textbox(s, Inches(0.9), Inches(5.4), Inches(6), Inches(0.3),
        'dokaz.net', size=12, color=BRAND_300, font='Consolas', bold=True)
textbox(s, Inches(0.9), Inches(5.75), Inches(9), Inches(0.3),
        'SOC 2 · ISO 27001 · HIPAA · cyber-insurance renewals',
        size=10, color=MUTED)


# =============================================================================
# Slide 2 — The risk
# =============================================================================
s = add_slide()
page_chrome(s, 'The risk',
            'Your backups are untested until something restores them.', 2, TOTAL)


def bullet_card(slide, left, top, w, h, tag, headline, body, accent=BRAND_400):
    rounded(slide, left, top, w, h, CARD_BG, line=STEEL_700)
    rect(slide, left, top + Inches(0.25), Inches(0.06), h - Inches(0.5), accent)
    textbox(slide, left + Inches(0.25), top + Inches(0.25), w - Inches(0.5), Inches(0.3),
            tag.upper(), size=9, bold=True, color=accent)
    textbox(slide, left + Inches(0.25), top + Inches(0.55), w - Inches(0.5), Inches(0.6),
            headline, size=15, bold=True, color=INK)
    textbox(slide, left + Inches(0.25), top + Inches(1.2), w - Inches(0.5), h - Inches(1.4),
            body, size=11, color=BODY)


bullet_card(s, Inches(0.6), Inches(2.4), Inches(6.05), Inches(2.35),
            'What you think you have',
            'A backup job that "succeeded" last night.',
            'A green tick on your S3 upload / cron job / RDS snapshot. '
            'It measures that a file was written — not whether the file, on restore, '
            'produces the data you actually expect.',
            accent=BRAND_400)

bullet_card(s, Inches(7.0), Inches(2.4), Inches(6.05), Inches(2.35),
            'What you actually have',
            'A dump you have never opened.',
            'Backup format changed after a Postgres upgrade. A column got dropped. '
            'An archive is truncated because the disk filled at 3am. You find out during the outage.',
            accent=DANGER)

bullet_card(s, Inches(0.6), Inches(4.95), Inches(6.05), Inches(2.15),
            'What your auditor asks for',
            'Restore-tested-in-the-last-12-months evidence.',
            'SOC 2 CC7.4 / A1.3, ISO 27001 A.8.13, and every cyber-insurance '
            'renewal from 2024 onward. Not the backup log — the restore result.',
            accent=PINK_400)

bullet_card(s, Inches(7.0), Inches(4.95), Inches(6.05), Inches(2.15),
            'What most teams do',
            'Screenshot a test restore, once a quarter.',
            'A manual test-restore is a person\'s afternoon. It is not repeatable, not '
            'signed, not verifiable. And it does not happen on the week the backup '
            'actually breaks.',
            accent=STEEL_500)


# =============================================================================
# Slide 3 — The fix
# =============================================================================
s = add_slide()
page_chrome(s, 'The fix',
            'A drill: a real restore, on a schedule, signed on the way out.',
            3, TOTAL)


def step_chip(slide, left, top, w, h, ordinal, title, body):
    rounded(slide, left, top, w, h, CARD_BG, line=STEEL_700)
    pill(slide, left + Inches(0.25), top + Inches(0.2), Inches(0.5), Inches(0.32),
         ordinal, bg=BRAND_500, fg=WHITE, size=11)
    textbox(slide, left + Inches(0.9), top + Inches(0.18), w - Inches(1.2), Inches(0.4),
            title, size=13, bold=True, color=INK)
    textbox(slide, left + Inches(0.3), top + Inches(0.75), w - Inches(0.6), h - Inches(0.9),
            body, size=10, color=BODY)


row_top = Inches(2.4)
row_h = Inches(1.7)
step_w = Inches(1.95)
step_gap = Inches(0.1)
left = Inches(0.6)
steps = [
    ('01', 'Provision', 'Isolated sandbox database. Ephemeral. Torn down when the drill ends.'),
    ('02', 'Fetch',     'Pull the dump you already produce. Hash it (SHA-256).'),
    ('03', 'Restore',   'pg_restore into the sandbox. A broken archive fails here — before your outage.'),
    ('04', 'Assert',    'Your data checks: row_count, table_exists, column_exists, no_nulls, SQL.'),
    ('05', 'Report',    'Render the PDF. Sign it Ed25519. Store the signature.'),
    ('06', 'Teardown',  'Destroy the sandbox and the working copy. Only the signed PDF survives.'),
]
for i, (ordinal, title, body) in enumerate(steps):
    step_chip(s, left + i * (step_w + step_gap), row_top, step_w, row_h, ordinal, title, body)

# Compact receipt mock underneath.
rounded(s, Inches(0.6), Inches(4.55), Inches(12.1), Inches(2.5), CARD_BG, line=STEEL_700)
textbox(s, Inches(0.9), Inches(4.75), Inches(6), Inches(0.4),
        'PROOF-OF-RECOVERY EVIDENCE', size=10, bold=True, color=STEEL_400)
pill(s, Inches(10.9), Inches(4.8), Inches(1.7), Inches(0.32), 'VERIFIED',
     bg=TEAL_500, fg=STEEL_950, size=10)

textbox(s, Inches(0.9), Inches(5.25), Inches(11), Inches(0.35),
        'database           production-primary',
        size=11, color=STEEL_200, font='Consolas')
textbox(s, Inches(0.9), Inches(5.6), Inches(11), Inches(0.35),
        'sha256(dump)       a1f9c2e4b7…d83f   matches',
        size=11, color=STEEL_200, font='Consolas')
textbox(s, Inches(0.9), Inches(5.95), Inches(11), Inches(0.35),
        'signature          ed25519:9f2c4b…a17b   valid',
        size=11, color=STEEL_200, font='Consolas')
textbox(s, Inches(0.9), Inches(6.3), Inches(11), Inches(0.35),
        'assertions passed  6 / 6      retain until 2033-05-22',
        size=11, color=STEEL_200, font='Consolas')

textbox(s, Inches(0.6), Inches(7.1), Inches(12), Inches(0.3),
        'Same six steps every drill. Same signed receipt every drill. Every signature verifies with the open-source dokaz-verify CLI — no one has to trust us.',
        size=10, color=MUTED, align=PP_ALIGN.CENTER)


# =============================================================================
# Slide 4 — What you get (product)
# =============================================================================
s = add_slide()
page_chrome(s, 'What you get',
            'Proof — not a green checkmark.', 4, TOTAL)


def feature_card(slide, left, top, w, h, title, body, accent=BRAND_400):
    rounded(slide, left, top, w, h, CARD_BG, line=STEEL_700)
    rect(slide, left, top + Inches(0.25), Inches(0.06), h - Inches(0.5), accent)
    textbox(slide, left + Inches(0.25), top + Inches(0.25), w - Inches(0.5), Inches(0.35),
            title, size=13, bold=True, color=INK)
    textbox(slide, left + Inches(0.25), top + Inches(0.7), w - Inches(0.5), h - Inches(0.85),
            body, size=10, color=BODY)


grid_left = Inches(0.6)
grid_top = Inches(2.4)
fw, fh = Inches(6.05), Inches(1.5)
fgap = Inches(0.15)

features = [
    ('Signed PDF evidence',       'Every drill produces an Ed25519-signed Proof-of-Recovery PDF, retained for 7 years by default. Independently verifiable — no need to trust us.', BRAND_400),
    ('Auditor share links',       'One-click tokenised URL. Your auditor sees the receipt, downloads the PDF, verifies the signature. No Dokaz account required. Revocable anytime.', PINK_400),
    ('Recovery-readiness score',  'Per-database A-F grade. Freshness of last passing drill × recent pass rate × latest outcome. The single glanceable answer to "how confident should I be?"', TEAL_400),
    ('Native alerting',           'Drill failed → Slack + PagerDuty + email + mobile push. Backup check-ins (heartbeats) fire the same paths when a scheduled job stops running.', BRAND_400),
    ('Public /verify page',       'Send an auditor the signed PDF + signature JSON. They open dokaz.net/verify, upload both files, and see "verified" against our published public key.', PINK_400),
    ('REST API + webhooks',       'Versioned /v1 with per-plan rate limits. Signed webhooks push drill.failed / heartbeat.down straight into your own systems.', TEAL_400),
]
for i, (title, body, accent) in enumerate(features):
    r, c = divmod(i, 2)
    feature_card(s, grid_left + c * (fw + fgap), grid_top + r * (fh + fgap), fw, fh, title, body, accent)


# =============================================================================
# Slide 5 — Pricing
# =============================================================================
s = add_slide()
page_chrome(s, 'Pricing',
            '$100 starter. Free to try, no card required. Daily drills on every tier.', 5, TOTAL)


def pricing_card(slide, left, top, w, h, name, price, period, cadence_label,
                 cadence_value, features, popular=False):
    border = PINK_400 if popular else STEEL_700
    line_w = Pt(2) if popular else Pt(0.75)
    rounded(slide, left, top, w, h, CARD_BG, line=border, line_w=line_w)
    if popular:
        pill(slide, left + Inches(0.4), top - Inches(0.15), Inches(2), Inches(0.3),
             'MOST POPULAR', bg=PINK_400, fg=STEEL_950, size=8)
    textbox(slide, left + Inches(0.3), top + Inches(0.35), w - Inches(0.6), Inches(0.35),
            name, size=14, bold=True, color=INK)
    textbox(slide, left + Inches(0.3), top + Inches(0.75), w - Inches(0.6), Inches(0.7),
            price, size=32, bold=True, color=BRAND_300)
    textbox(slide, left + Inches(0.3), top + Inches(0.95), w - Inches(0.6), Inches(0.4),
            period, size=11, color=STEEL_400)
    rounded(slide, left + Inches(0.3), top + Inches(1.65), w - Inches(0.6), Inches(0.5),
            STEEL_800, line=STEEL_700)
    textbox(slide, left + Inches(0.3), top + Inches(1.68), w - Inches(0.6), Inches(0.22),
            cadence_label.upper(), size=8, bold=True, color=STEEL_400, align=PP_ALIGN.CENTER)
    textbox(slide, left + Inches(0.3), top + Inches(1.88), w - Inches(0.6), Inches(0.28),
            cadence_value, size=12, bold=True, color=BRAND_300, align=PP_ALIGN.CENTER)
    y = top + Inches(2.35)
    for f in features:
        textbox(slide, left + Inches(0.3), y, w - Inches(0.6), Inches(0.3),
                '· ' + f, size=9.5, color=BODY)
        y += Inches(0.32)


card_w = Inches(3.0)
card_h = Inches(4.7)
card_gap = Inches(0.15)
card_top = Inches(2.4)
left0 = Inches(0.6)

pricing_card(s, left0, card_top, card_w, card_h,
             'Starter', '$100', '/mo', 'Capacity', '5 databases',
             ['Daily drills · signed PDF · 7-year retention',
              'Unlimited seats · 25 backup check-ins',
              'Auditor share links', 'Slack + PagerDuty', 'Email support'])

pricing_card(s, left0 + (card_w + card_gap), card_top, card_w, card_h,
             'Growth', '$300', '/mo', 'Capacity', '25 databases',
             ['Daily drills · everything in Starter',
              '50 backup check-ins',
              'Unlimited API keys · signed webhooks',
              'Priority support'],
             popular=True)

pricing_card(s, left0 + (card_w + card_gap) * 2, card_top, card_w, card_h,
             'Grounded', '$600', '/mo', 'Capacity', '100 databases',
             ['Daily drills · everything in Growth',
              'Unlimited backup check-ins',
              'Priority + dedicated channel', 'SSO (roadmap)'])

pricing_card(s, left0 + (card_w + card_gap) * 3, card_top, card_w, card_h,
             'Enterprise', 'Custom', '', 'Regulated / hands-on', 'Sales-led',
             ['BYO runner in your VPC · SSO/SAML',
              'BAA / DPA · hourly cadence',
              'Auditor read-only accounts',
              'Custom SLA · named account manager'])

textbox(s, Inches(0.6), Inches(7.2), Inches(12), Inches(0.3),
        'Free trial, no card required — one real database at weekly cadence. '
        'Every paid tier includes daily drills, unlimited seats, 7-year retention.',
        size=10, color=MUTED, align=PP_ALIGN.CENTER)


# =============================================================================
# Slide 6 — Send us an auditor's test link right now
# =============================================================================
s = add_slide()
page_chrome(s, 'Try it now',
            'Verify a signed report before you sign up.', 6, TOTAL)

# Two-column CTA layout: (left) the /verify page URL big; (right) how it works.
rounded(s, Inches(0.6), Inches(2.4), Inches(6.05), Inches(4.7),
        CARD_BG, line=BRAND_400, line_w=Pt(2))
textbox(s, Inches(0.9), Inches(2.7), Inches(5.4), Inches(0.35),
        'STEP 1', size=9, bold=True, color=PINK_300)
textbox(s, Inches(0.9), Inches(3.0), Inches(5.4), Inches(0.6),
        'Download a real signed report.',
        size=18, bold=True, color=INK)
textbox(s, Inches(0.9), Inches(3.7), Inches(5.4), Inches(1.5),
        'We keep a sample Proof-of-Recovery PDF and its signature at:\n\n'
        'dokaz.net/docs (Evidence · sample)\n\n'
        'The PDF is a real drill run against our sample fixture. The signature JSON '
        'is our production Ed25519 signing key.',
        size=11, color=BODY)
pill(s, Inches(0.9), Inches(6.35), Inches(3), Inches(0.5),
     'dokaz.net/docs', bg=BRAND_500, fg=WHITE, size=12)

rounded(s, Inches(7.0), Inches(2.4), Inches(6.05), Inches(4.7),
        CARD_BG, line=TEAL_400, line_w=Pt(2))
textbox(s, Inches(7.3), Inches(2.7), Inches(5.4), Inches(0.35),
        'STEP 2', size=9, bold=True, color=TEAL_300)
textbox(s, Inches(7.3), Inches(3.0), Inches(5.4), Inches(0.6),
        'Verify it in your browser.',
        size=18, bold=True, color=INK)
textbox(s, Inches(7.3), Inches(3.7), Inches(5.4), Inches(1.5),
        'Open dokaz.net/verify. Upload the PDF + signature JSON. The page checks '
        'the signature against our published public key and returns a receipt.\n\n'
        'No account, no cookies, no data stored. If it says verified, the report is '
        'authentic and unaltered.',
        size=11, color=BODY)
pill(s, Inches(7.3), Inches(6.35), Inches(3), Inches(0.5),
     'dokaz.net/verify', bg=TEAL_500, fg=STEEL_950, size=12)


# =============================================================================
# Slide 7 — Sign up
# =============================================================================
s = add_slide()

glow = rect(s, Inches(9), Inches(-2), Inches(6), Inches(6), BRAND_900)
glow.line.fill.background()

textbox(s, Inches(0.9), Inches(2.4), Inches(11), Inches(0.9),
        'You already write backups.', size=42, bold=True, color=INK)
textbox(s, Inches(0.9), Inches(3.15), Inches(11), Inches(0.9),
        'Let\'s prove they restore.', size=42, bold=True, color=PINK_300)

textbox(s, Inches(0.9), Inches(4.5), Inches(11), Inches(1.6),
        'Sign up at dokaz.net. Free trial, no card required. Run a drill against our sample dataset in 30 seconds — '
        'you\'ll see the exact signed PDF your auditor will see. Then connect your own '
        'pg_dump and prove it on your real backup, still inside the trial.',
        size=15, color=BODY)

pill(s, Inches(0.9), Inches(6.2), Inches(3.5), Inches(0.6),
     'dokaz.net/signup', bg=BRAND_500, fg=WHITE, size=16)
pill(s, Inches(4.7), Inches(6.2), Inches(3.5), Inches(0.6),
     'Book a 20-min demo', bg=STEEL_800, fg=INK, size=13)

textbox(s, Inches(0.9), Inches(7.05), Inches(11), Inches(0.35),
        'Questions: sales@dokaz.net · Verify a report: dokaz.net/verify · Status: dokaz.net/status',
        size=10, color=MUTED)


# ---- save ----
out = Path(__file__).parent / 'Dokaz_Sales_Deck.pptx'
prs.save(str(out))
print(f'Wrote {out}')
