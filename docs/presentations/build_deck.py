"""Generate Dokaz_Investor_Deck.pptx — a 10-slide investor brief.

Palette matches the shipped Dokaz brand (steel + baby-royal-blue + pink +
sea-teal — see tailwind.config.js and branding/design-system/). Voice matches
the marketing site: direct, technical, unflinching. Content reflects what
dokaz.net actually is today (Ed25519-signed evidence, auditor share links,
readiness score, Slack/PagerDuty alerts, /verify, $99/$299/$799/Enterprise).

Run: `python docs/presentations/build_deck.py` — writes Dokaz_Investor_Deck.pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# =============================================================================
# Brand palette — mirrors branding/design-system/tokens/colors.css.
# Dokaz is a dark-only brand; the deck is dark to match how the app renders.
# =============================================================================
BRAND_900 = RGBColor(0x24, 0x3A, 0x78)
BRAND_700 = RGBColor(0x3F, 0x63, 0xC4)
BRAND_500 = RGBColor(0x5B, 0x8D, 0xEF)
BRAND_400 = RGBColor(0x6E, 0x9B, 0xF0)
BRAND_300 = RGBColor(0x8C, 0xB1, 0xF7)
BRAND_100 = RGBColor(0xDB, 0xE7, 0xFF)

PINK_500  = RGBColor(0xEC, 0x48, 0x99)
PINK_400  = RGBColor(0xF4, 0x8F, 0xB1)
PINK_300  = RGBColor(0xF9, 0xA8, 0xD4)

TEAL_500  = RGBColor(0x2B, 0xA8, 0x88)
TEAL_400  = RGBColor(0x56, 0xC5, 0x96)
TEAL_300  = RGBColor(0x5F, 0xCB, 0xAC)

DANGER    = RGBColor(0xFF, 0x6B, 0x81)
AMBER     = RGBColor(0xF0, 0xB4, 0x29)

STEEL_950 = RGBColor(0x11, 0x15, 0x1A)
STEEL_900 = RGBColor(0x1B, 0x20, 0x27)
STEEL_800 = RGBColor(0x25, 0x2C, 0x35)
STEEL_700 = RGBColor(0x33, 0x40, 0x5A)
STEEL_600 = RGBColor(0x47, 0x54, 0x67)
STEEL_500 = RGBColor(0x5D, 0x6C, 0x80)
STEEL_400 = RGBColor(0x7A, 0x8A, 0x9C)
STEEL_300 = RGBColor(0x9A, 0xA7, 0xB6)
STEEL_200 = RGBColor(0xC4, 0xCF, 0xDB)
STEEL_100 = RGBColor(0xE5, 0xEC, 0xF2)
STEEL_50  = RGBColor(0xF4, 0xF7, 0xFA)

WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0xED, 0xF1, 0xF6)
BODY      = RGBColor(0xD7, 0xDE, 0xE7)
MUTED     = RGBColor(0x9A, 0xA7, 0xB6)
FAINT     = RGBColor(0x6B, 0x7A, 0x8D)

# Canonical card surface (matches .card in the app).
CARD_BG   = RGBColor(0x2E, 0x36, 0x40)

# =============================================================================
# Deck primitives
# =============================================================================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

TOTAL = 10


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = STEEL_900
    bg.shadow.inherit = False
    return s


def textbox(slide, left, top, width, height, text, *,
            size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def rounded(slide, left, top, width, height, fill, line=None, line_w=Pt(0.75)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.10
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_w
    shape.shadow.inherit = False
    return shape


def pill(slide, left, top, width, height, text, *, bg=BRAND_500, fg=WHITE, size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    shape.shadow.inherit = False
    tb = shape.text_frame
    tb.margin_left = tb.margin_right = Inches(0.05)
    tb.margin_top = tb.margin_bottom = 0
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    return shape


def page_chrome(slide, eyebrow, title, slide_num, total):
    textbox(slide, Inches(0.6), Inches(0.35), Inches(2), Inches(0.35),
            'Dokaz', size=13, bold=True, color=INK)
    textbox(slide, Inches(12.4), Inches(0.35), Inches(0.6), Inches(0.35),
            f'{slide_num:02d} / {total:02d}',
            size=9, color=STEEL_500, align=PP_ALIGN.RIGHT, font='Consolas')
    textbox(slide, Inches(0.6), Inches(0.85), Inches(6), Inches(0.35),
            eyebrow.upper(), size=10, bold=True, color=PINK_300)
    textbox(slide, Inches(0.6), Inches(1.15), Inches(12), Inches(0.9),
            title, size=30, bold=True, color=INK)
    rect(slide, Inches(0.6), Inches(2.05), Inches(0.6), Inches(0.04), PINK_400)


# =============================================================================
# Slide 1 — Cover
# =============================================================================
s = add_slide()

glow_l = rect(s, Inches(-2), Inches(-1), Inches(6), Inches(5), BRAND_900)
glow_l.line.fill.background()
glow_r = rect(s, Inches(9.5), Inches(4.5), Inches(6), Inches(4), STEEL_800)
glow_r.line.fill.background()

textbox(s, Inches(0.9), Inches(2.5), Inches(6), Inches(0.9),
        'Dokaz', size=68, bold=True, color=WHITE)
textbox(s, Inches(0.9), Inches(3.5), Inches(9), Inches(0.5),
        'Backup verification you can independently prove.',
        size=22, color=PINK_300)
textbox(s, Inches(0.9), Inches(4.1), Inches(11), Inches(0.9),
        'Dokaz periodically restores your database backups in an isolated sandbox, '
        'asserts the data is intact, and emits a signed Proof-of-Recovery PDF '
        'any auditor can verify without trusting us.',
        size=13, color=BODY)

textbox(s, Inches(0.9), Inches(5.4), Inches(6), Inches(0.3),
        'dokaz.net', size=12, color=BRAND_300, font='Consolas', bold=True)
textbox(s, Inches(0.9), Inches(5.75), Inches(9), Inches(0.3),
        'Series Seed · B2B infrastructure · Postgres today, MySQL / MongoDB roadmap',
        size=10, color=MUTED)

textbox(s, Inches(9.3), Inches(6.9), Inches(3.6), Inches(0.35),
        'Investor briefing — 2026',
        size=9, color=STEEL_500, align=PP_ALIGN.RIGHT, font='Consolas')


# =============================================================================
# Slide 2 — The problem
# =============================================================================
s = add_slide()
page_chrome(s, 'The problem',
            'Backups fail silently. You find out during the outage.', 2, TOTAL)


def stat_card(slide, left, top, w, h, label, value, sub, accent=BRAND_400):
    rounded(slide, left, top, w, h, CARD_BG, line=STEEL_700)
    textbox(slide, left + Inches(0.3), top + Inches(0.25), w - Inches(0.6), Inches(0.35),
            label.upper(), size=10, bold=True, color=STEEL_400)
    textbox(slide, left + Inches(0.3), top + Inches(0.6), w - Inches(0.6), Inches(0.9),
            value, size=32, bold=True, color=accent)
    textbox(slide, left + Inches(0.3), top + Inches(1.55), w - Inches(0.6), Inches(0.9),
            sub, size=10, color=BODY)


card_w, card_h = Inches(3.85), Inches(2.55)
card_gap = Inches(0.25)
card_top = Inches(2.4)
left0 = Inches(0.6)

stat_card(s, left0, card_top, card_w, card_h,
          'Cyber-insurance renewal', 'Every 12 mo',
          'Carriers ask for "restore-tested-in-last-12-months" evidence. '
          '"Our backup job succeeded" does not satisfy it.',
          accent=BRAND_400)

stat_card(s, left0 + (card_w + card_gap), card_top, card_w, card_h,
          'SOC 2 CC7.4 / A1.3', 'Restore evidence',
          'Auditors expect a *restored* backup with data checks — not a green '
          'checkmark on a job log. Most teams fake this quarterly.',
          accent=PINK_400)

stat_card(s, left0 + (card_w + card_gap) * 2, card_top, card_w, card_h,
          'Actual outages', 'First-restore',
          'Most companies discover their backup format changed / a column '
          'got dropped / the archive is truncated the day they need it.',
          accent=DANGER)

textbox(s, Inches(0.6), Inches(5.4), Inches(12), Inches(1.6),
        '"Successful backup job" measures that a file was written. '
        'It does not measure whether the file, on restore, produces the data you '
        'expect. Every regulated B2B company already writes backups — very few '
        'actually verify them, because verifying means running a real restore '
        'somewhere. Dokaz is that somewhere.',
        size=14, color=BODY)


# =============================================================================
# Slide 3 — The solution
# =============================================================================
s = add_slide()
page_chrome(s, 'The solution',
            'A drill is a real restore, on a schedule, cryptographically signed.',
            3, TOTAL)


def step_chip(slide, left, top, w, h, ordinal, title, body):
    rounded(slide, left, top, w, h, CARD_BG, line=STEEL_700)
    pill(slide, left + Inches(0.25), top + Inches(0.2), Inches(0.5), Inches(0.32),
         ordinal, bg=BRAND_500, fg=WHITE, size=11)
    textbox(slide, left + Inches(0.9), top + Inches(0.18), w - Inches(1.2), Inches(0.4),
            title, size=13, bold=True, color=INK)
    textbox(slide, left + Inches(0.3), top + Inches(0.75), w - Inches(0.6), h - Inches(0.9),
            body, size=10, color=BODY)


row_top = Inches(2.4)
row_h = Inches(1.7)
step_w = Inches(1.95)
step_gap = Inches(0.1)
left = Inches(0.6)
steps = [
    ('01', 'Provision', 'Spin up an isolated sandbox database. Ephemeral, torn down after the drill.'),
    ('02', 'Fetch',     'Pull the customer\'s dump. Hash it (SHA-256) — that hash goes in the receipt.'),
    ('03', 'Restore',   'pg_restore into the sandbox. If the archive is corrupt, we find out here.'),
    ('04', 'Assert',    'Run the customer\'s data checks: row_count, table_exists, no_nulls, SQL.'),
    ('05', 'Report',    'Render the PDF. Sign it with our Ed25519 key. Store the signature.'),
    ('06', 'Teardown',  'Destroy the sandbox and the working copy. The evidence PDF survives.'),
]
for i, (ordinal, title, body) in enumerate(steps):
    step_chip(s, left + i * (step_w + step_gap), row_top, step_w, row_h, ordinal, title, body)

rounded(s, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.5), CARD_BG, line=STEEL_700)
textbox(s, Inches(0.9), Inches(4.8), Inches(6), Inches(0.4),
        'PROOF-OF-RECOVERY EVIDENCE (SAMPLE OUTPUT)',
        size=10, bold=True, color=STEEL_400)
pill(s, Inches(10.9), Inches(4.85), Inches(1.7), Inches(0.32), 'VERIFIED',
     bg=TEAL_500, fg=STEEL_950, size=10)

textbox(s, Inches(0.9), Inches(5.3), Inches(11), Inches(0.35),
        'sha256(dump)       a1f9c2e4b7…d83f   matches',
        size=11, color=STEEL_200, font='Consolas')
textbox(s, Inches(0.9), Inches(5.65), Inches(11), Inches(0.35),
        'signature          ed25519:9f2c4b…a17b   valid',
        size=11, color=STEEL_200, font='Consolas')
textbox(s, Inches(0.9), Inches(6.0), Inches(11), Inches(0.35),
        'assertions passed  6 / 6',
        size=11, color=STEEL_200, font='Consolas')
textbox(s, Inches(0.9), Inches(6.35), Inches(11), Inches(0.35),
        'retain until       2033-05-22',
        size=11, color=STEEL_200, font='Consolas')

textbox(s, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
        'The moat: every signature is verifiable with our open-source dokaz-verify CLI against the '
        'public key at dokaz.net/.well-known/evidence-signing-keys.pem — auditors do not need to trust Dokaz.',
        size=10, color=MUTED, align=PP_ALIGN.CENTER)


# =============================================================================
# Slide 4 — Onboarding
# =============================================================================
s = add_slide()
page_chrome(s, 'Onboarding',
            'From signup to a signed report in 30 seconds.', 4, TOTAL)

row_top = Inches(2.6)
row_h = Inches(2.6)
step_w = Inches(2.95)
step_gap = Inches(0.15)
left = Inches(0.6)

onboarding_steps = [
    ('1', 'Sign up',
     'Email + password. $1 first month; card required, cancel anytime. '
     'The trial gives one real database at weekly cadence, so the buyer proves the product on their own dump.'),
    ('2', 'Run the sample',
     'A one-click drill against our built-in Postgres fixture runs the full '
     'six-step pipeline in ~2 seconds and produces a real signed PDF. This is the '
     'aha moment — before they connect anything.'),
    ('3', 'Connect a database',
     'Point Dokaz at a pg_dump. We fetch, restore in an isolated sandbox, '
     'assert row counts / table exists / no nulls, and sign the result. Weekly on '
     'Starter and Growth, daily on Scale.'),
    ('4', 'Share the evidence',
     'One tokenised URL for the auditor. Read-only, no Dokaz account, expires in 30 days. '
     'Signature verifies against the same public key the CLI uses.'),
]
for i, (ordinal, title, body) in enumerate(onboarding_steps):
    step_chip(s, left + i * (step_w + step_gap), row_top, step_w, row_h, ordinal, title, body)

textbox(s, Inches(0.6), Inches(5.5), Inches(12), Inches(1.6),
        'The 30-second sample drill is the funnel. It short-circuits the classic '
        '"free trial that never actually shows the value" pattern — the prospect sees '
        'the exact signed PDF their auditor will see, before they touch their own '
        'infra. First-time-value is the sale.',
        size=13, color=BODY)


# =============================================================================
# Slide 5 — Product (what we have built)
# =============================================================================
s = add_slide()
page_chrome(s, 'Product',
            'The evidence chain, and the surfaces around it.', 5, TOTAL)


def feature_card(slide, left, top, w, h, title, body, accent=BRAND_400):
    rounded(slide, left, top, w, h, CARD_BG, line=STEEL_700)
    rect(slide, left, top + Inches(0.25), Inches(0.06), h - Inches(0.5), accent)
    textbox(slide, left + Inches(0.25), top + Inches(0.25), w - Inches(0.5), Inches(0.35),
            title, size=12, bold=True, color=INK)
    textbox(slide, left + Inches(0.25), top + Inches(0.65), w - Inches(0.5), h - Inches(0.8),
            body, size=9.5, color=BODY)


grid_left = Inches(0.6)
grid_top = Inches(2.35)
fw, fh = Inches(4.05), Inches(1.5)
fgap = Inches(0.1)

features = [
    ('Six-step drill pipeline',   'Provision → fetch → restore → assert → report → teardown. '
                                  'Every step has structured logs, a status, and a 16 KiB output snippet.', BRAND_400),
    ('Ed25519-signed PDF',        'Detached signature over the SHA-256 of the PDF. Retention configurable to 7 years. '
                                  'Signing key rotation preserves verifiability of old reports.', PINK_400),
    ('Open-source verifier',      'dokaz-verify CLI is 200 lines of Go. Zero Dokaz-specific runtime — auditors run it on their machine '
                                  'against the published public key.', TEAL_400),
    ('Auditor share links',       'Tokenised, expiring, revocable URLs. Read-only receipt + downloadable PDF + '
                                  '"verify independently" affordance. No Dokaz account required.', PINK_400),
    ('Recovery-readiness score',  'Per-database A-F grade fusing freshness, pass rate, and last outcome. '
                                  'The single glanceable answer to "how confident should I be if I had to restore now?"', BRAND_400),
    ('Native alerting',           'Drill fails → Slack + PagerDuty + email + mobile push. Backup check-ins '
                                  '(heartbeats) fire the same paths when a job stops running.', TEAL_400),
    ('Public /verify page',       'An auditor pastes a PDF + signature JSON, the browser checks it against '
                                  'the published key. Zero trust in Dokaz. Renders as the same receipt.', BRAND_400),
    ('Evidence bundle',           'One-click annual ZIP: every signed PDF + signatures + a manifest CSV. '
                                  'The exact artifact the auditor asks for at renewal.', PINK_400),
    ('REST API + webhooks',       'Versioned /v1 with per-plan rate limits. Signed webhooks for drill.failed / '
                                  'heartbeat.down. Powers the responder mobile app.', TEAL_400),
]
for i, (title, body, accent) in enumerate(features):
    r, c = divmod(i, 3)
    feature_card(s, grid_left + c * (fw + fgap), grid_top + r * (fh + fgap), fw, fh, title, body, accent)


# =============================================================================
# Slide 6 — Comparison (a failed drill vs a successful one)
# =============================================================================
s = add_slide()
page_chrome(s, 'Comparison',
            'A failed drill and a successful one — both make the same evidence.',
            6, TOTAL)


def drill_card(slide, left, top, w, h, *, verdict, verdict_color, target,
               drill_id, restore_time, rows, assertions, evidence_label,
               step_states, callout_msg=None, callout_color=None):
    rounded(slide, left, top, w, h, CARD_BG, line=STEEL_700)

    textbox(slide, left + Inches(0.3), top + Inches(0.25), w - Inches(2.5), Inches(0.4),
            target, size=15, bold=True, color=INK)
    pill(slide, left + w - Inches(1.5), top + Inches(0.3), Inches(1.2), Inches(0.32),
         verdict, bg=verdict_color, fg=STEEL_950, size=10)
    textbox(slide, left + Inches(0.3), top + Inches(0.7), w - Inches(0.6), Inches(0.3),
            f'drill_{drill_id} · pg_dump -Fc',
            size=9, color=STEEL_400, font='Consolas')

    y = top + Inches(1.1)
    if callout_msg:
        rounded(slide, left + Inches(0.3), y, w - Inches(0.6), Inches(0.42),
                CARD_BG, line=callout_color)
        textbox(slide, left + Inches(0.45), y + Inches(0.05), w - Inches(0.9), Inches(0.35),
                callout_msg, size=9, color=callout_color, font='Consolas')
        y = y + Inches(0.55)

    tile_top = y + Inches(0.05)
    tile_w = (w - Inches(0.9)) / 4
    tile_gap = Inches(0.05)
    for i, (label, value) in enumerate(
            [('Restore time', restore_time), ('Rows verified', rows),
             ('Assertions', assertions), ('Evidence', evidence_label)]):
        tx = left + Inches(0.3) + i * (tile_w + tile_gap)
        rounded(slide, tx, tile_top, tile_w, Inches(0.75), STEEL_800, line=STEEL_700)
        textbox(slide, tx + Inches(0.1), tile_top + Inches(0.05), tile_w - Inches(0.2), Inches(0.25),
                label.upper(), size=7, bold=True, color=STEEL_400)
        textbox(slide, tx + Inches(0.1), tile_top + Inches(0.3), tile_w - Inches(0.2), Inches(0.4),
                value, size=12, bold=True, color=INK)

    chip_top = tile_top + Inches(0.95)
    chip_names = ['provision', 'fetch', 'restore', 'assert', 'report', 'teardown']
    chip_w = Inches(0.85)
    for i, (name, state) in enumerate(zip(chip_names, step_states)):
        bg = TEAL_500 if state == 'passed' else DANGER if state == 'failed' \
            else STEEL_600 if state in ('skipped', 'pending') else BRAND_500
        pill(slide, left + Inches(0.3) + i * (chip_w + Inches(0.05)), chip_top,
             chip_w, Inches(0.28), name, bg=bg, fg=STEEL_950, size=8)


drill_card(s, Inches(0.6), Inches(2.4), Inches(6.05), Inches(4.7),
           verdict='PASSED', verdict_color=TEAL_500,
           target='production-primary', drill_id='a1f9c2e4',
           restore_time='2m 14s', rows='4.2M', assertions='6 / 6',
           evidence_label='Signed',
           step_states=['passed'] * 6)

drill_card(s, Inches(7.0), Inches(2.4), Inches(6.05), Inches(4.7),
           verdict='FAILED', verdict_color=DANGER,
           target='billing-primary', drill_id='f04ab991',
           restore_time='1m 02s', rows='—', assertions='3 / 4',
           evidence_label='Withheld',
           callout_msg='assert no_nulls: column user_id had 3 nulls',
           callout_color=DANGER,
           step_states=['passed', 'passed', 'passed', 'failed', 'skipped', 'passed'])

textbox(s, Inches(0.6), Inches(7.15), Inches(12), Inches(0.3),
        'Same pipeline, same receipt shape, opposite verdicts. Both are auditor-usable evidence — one proves recoverability, one proves the failure was discovered *before* the customer needed the restore.',
        size=10, color=MUTED, align=PP_ALIGN.CENTER)


# =============================================================================
# Slide 7 — Market
# =============================================================================
s = add_slide()
page_chrome(s, 'Market',
            'Every regulated B2B company with a database is a buyer.', 7, TOTAL)


def market_col(slide, left, top, w, h, header, big, sub, bullets):
    rounded(slide, left, top, w, h, CARD_BG, line=STEEL_700)
    textbox(slide, left + Inches(0.3), top + Inches(0.25), w - Inches(0.6), Inches(0.35),
            header.upper(), size=10, bold=True, color=STEEL_400)
    textbox(slide, left + Inches(0.3), top + Inches(0.6), w - Inches(0.6), Inches(0.8),
            big, size=28, bold=True, color=BRAND_300)
    textbox(slide, left + Inches(0.3), top + Inches(1.35), w - Inches(0.6), Inches(0.4),
            sub, size=10, color=MUTED)
    y = top + Inches(1.85)
    for b in bullets:
        textbox(slide, left + Inches(0.3), y, w - Inches(0.6), Inches(0.3),
                '· ' + b, size=10, color=BODY)
        y += Inches(0.3)


market_col(s, Inches(0.6), Inches(2.4), Inches(3.9), Inches(4.6),
           'Immediate wedge — SOC 2 in flight', '~50k',
           'US SMB SaaS at Series A-C. Auditor at the door, backup-restore evidence '
           'in the checklist, no dedicated infra headcount.',
           ['Vanta / Drata already sold', 'Buys pain, not vision',
            'CTO is champion + buyer'])

market_col(s, Inches(4.7), Inches(2.4), Inches(3.9), Inches(4.6),
           'Adjacent — cyber-insurance renewals', '$8.4B',
           'US cyber-insurance policies annually. Every renewal now asks for '
           'restore-tested-in-last-12-months evidence. Nothing else on the market emits it.',
           ['Insurance broker referrals',
            'Broker becomes distribution',
            'Buys once a year, forever'])

market_col(s, Inches(8.8), Inches(2.4), Inches(3.9), Inches(4.6),
           'Enterprise — HIPAA / ISO 27001 / DORA', '5x TAM',
           'Healthcare, EU financial services (DORA 2025), regulated fintech. '
           'Same evidence chain, tighter retention, VPC-hosted runner.',
           ['Longer cycle, bigger ACV',
            'SSO / DPA / BAA required',
            'Follows the SMB wedge in year 2'])


# =============================================================================
# Slide 8 — Pricing
# =============================================================================
s = add_slide()
page_chrome(s, 'Pricing',
            'Four tiers. Priced on drill volume — the marginal cost.', 8, TOTAL)


def pricing_card(slide, left, top, w, h, name, price, period, cadence_label,
                 cadence_value, features, popular=False):
    border = PINK_400 if popular else STEEL_700
    line_w = Pt(2) if popular else Pt(0.75)
    rounded(slide, left, top, w, h, CARD_BG, line=border, line_w=line_w)
    if popular:
        pill(slide, left + Inches(0.4), top - Inches(0.15), Inches(2), Inches(0.3),
             'MOST POPULAR', bg=PINK_400, fg=STEEL_950, size=8)
    textbox(slide, left + Inches(0.3), top + Inches(0.35), w - Inches(0.6), Inches(0.35),
            name, size=14, bold=True, color=INK)
    textbox(slide, left + Inches(0.3), top + Inches(0.75), w - Inches(0.6), Inches(0.7),
            price, size=32, bold=True, color=BRAND_300)
    textbox(slide, left + Inches(0.3), top + Inches(0.95), w - Inches(0.6), Inches(0.4),
            period, size=11, color=STEEL_400)
    rounded(slide, left + Inches(0.3), top + Inches(1.65), w - Inches(0.6), Inches(0.5),
            STEEL_800, line=STEEL_700)
    textbox(slide, left + Inches(0.3), top + Inches(1.68), w - Inches(0.6), Inches(0.22),
            cadence_label.upper(), size=8, bold=True, color=STEEL_400, align=PP_ALIGN.CENTER)
    textbox(slide, left + Inches(0.3), top + Inches(1.88), w - Inches(0.6), Inches(0.28),
            cadence_value, size=12, bold=True, color=BRAND_300, align=PP_ALIGN.CENTER)
    y = top + Inches(2.35)
    for f in features:
        textbox(slide, left + Inches(0.3), y, w - Inches(0.6), Inches(0.3),
                '· ' + f, size=9.5, color=BODY)
        y += Inches(0.32)


card_w = Inches(3.0)
card_h = Inches(4.7)
card_gap = Inches(0.15)
card_top = Inches(2.4)
left0 = Inches(0.6)

pricing_card(s, left0, card_top, card_w, card_h,
             'Starter', '$99', '/mo', 'Drill frequency', 'Weekly',
             ['5 databases', '3 seats', 'Signed PDF · 7-year retention',
              'Auditor share links', 'Slack + PagerDuty', 'Email support'])

pricing_card(s, left0 + (card_w + card_gap), card_top, card_w, card_h,
             'Growth', '$299', '/mo', 'Drill frequency', 'Weekly',
             ['25 databases', '10 seats', 'Everything in Starter',
              'API access · signed webhooks', 'Priority support'],
             popular=True)

pricing_card(s, left0 + (card_w + card_gap) * 2, card_top, card_w, card_h,
             'Scale', '$799', '/mo', 'Drill frequency', 'Daily',
             ['Unlimited databases', 'Unlimited seats',
              'Everything in Growth', 'Priority + dedicated channel', 'SSO (roadmap)'])

pricing_card(s, left0 + (card_w + card_gap) * 3, card_top, card_w, card_h,
             'Enterprise', 'From $1.5k', '/mo', 'Drill frequency', 'Custom',
             ['Hourly + custom cadence', 'Auditor read-only accounts',
              'BYO runner in your VPC', 'Longer retention · custom SLA',
              'Named account manager'])

textbox(s, Inches(0.6), Inches(7.2), Inches(12), Inches(0.3),
        'Trial: $1 first month, one real database at weekly cadence — the buyer '
        'proves the product on their own dump before the paywall re-arms.',
        size=10, color=MUTED, align=PP_ALIGN.CENTER)


# =============================================================================
# Slide 9 — Status (where we are today)
# =============================================================================
s = add_slide()
page_chrome(s, 'Status', 'Where we stand today.', 9, TOTAL)


def status_card(slide, left, top, w, h, header, big, sub):
    rounded(slide, left, top, w, h, CARD_BG, line=STEEL_700)
    textbox(slide, left + Inches(0.3), top + Inches(0.25), w - Inches(0.6), Inches(0.35),
            header.upper(), size=10, bold=True, color=STEEL_400)
    textbox(slide, left + Inches(0.3), top + Inches(0.65), w - Inches(0.6), Inches(1),
            big, size=22, bold=True, color=TEAL_300)
    textbox(slide, left + Inches(0.3), top + Inches(1.7), w - Inches(0.6), h - Inches(2),
            sub, size=10, color=BODY)


card_w = Inches(3.9)
card_h = Inches(2.1)
card_gap = Inches(0.15)
card_top = Inches(2.4)
left0 = Inches(0.6)

status_card(s, left0, card_top, card_w, card_h,
            'Product', 'Live at dokaz.net',
            'Full evidence chain shipping. Public verifier, auditor share links, readiness score, '
            'native Slack + PagerDuty, mobile responder app.')
status_card(s, left0 + (card_w + card_gap), card_top, card_w, card_h,
            'Infrastructure', 'Self-migrating',
            'Fly.io single-machine + Neon Postgres. Migrations run at server startup — deploys '
            'are one reliable step. Evidence keys persisted as Fly secrets.')
status_card(s, left0 + 2 * (card_w + card_gap), card_top, card_w, card_h,
            'Security posture', 'Argon2id + Ed25519',
            'Passwords Argon2id, sessions HTTPS-only, evidence signed Ed25519 with rotation-tolerant '
            'verification, at-rest AES-GCM with per-account keys (crypto-shred).')

card_top2 = Inches(4.65)
status_card(s, left0, card_top2, card_w, card_h,
            'Compliance chain', 'Verifiable',
            'Every PDF signature verifies against the public key at '
            'dokaz.net/.well-known/evidence-signing-keys.pem via the open-source dokaz-verify CLI.')
status_card(s, left0 + (card_w + card_gap), card_top2, card_w, card_h,
            'Team', 'One founder-engineer',
            'Sole builder to date. All backend, frontend, mobile, brand, and this deck. '
            'First-hire slot at $1M ARR: security-focused eng #2.')
status_card(s, left0 + 2 * (card_w + card_gap), card_top2, card_w, card_h,
            'Traction', 'Pre-launch',
            'Marketing live, pipeline pre-outbound. Design system + brand shipped '
            '2026 Q3. Fundraise timed to first-customer conversion.')


# =============================================================================
# Slide 10 — Plan (next 90 days)
# =============================================================================
s = add_slide()
page_chrome(s, 'Plan', 'The next 90 days.', 10, TOTAL)


def timeline_card(slide, left, top, w, h, when, title, items, color):
    rounded(slide, left, top, w, h, CARD_BG, line=STEEL_700)
    rect(slide, left, top + Inches(0.25), Inches(0.06), h - Inches(0.5), color)
    textbox(slide, left + Inches(0.25), top + Inches(0.25), w - Inches(0.5), Inches(0.3),
            when.upper(), size=9, bold=True, color=color)
    textbox(slide, left + Inches(0.25), top + Inches(0.55), w - Inches(0.5), Inches(0.45),
            title, size=15, bold=True, color=INK)
    y = top + Inches(1.15)
    for item in items:
        textbox(slide, left + Inches(0.25), y, w - Inches(0.5), Inches(0.35),
                '· ' + item, size=10, color=BODY)
        y += Inches(0.32)


card_w = Inches(3.9)
card_h = Inches(4.7)
card_gap = Inches(0.15)
card_top = Inches(2.4)
left0 = Inches(0.6)

timeline_card(s, left0, card_top, card_w, card_h,
              'Days 0-30', 'Outbound + first ten customers', [
                  'ICP: US SMB SaaS Series A-C, SOC 2 in flight',
                  '3-5 cold-email variants keyed on segment',
                  'Enrichment: SOC 2 badge, hiring signal, insurance renewal window',
                  'Weekly demo cadence on Loom, sub-3-minute',
                  'Goal: 10 paid pilots at $299 (Growth)',
              ], BRAND_400)
timeline_card(s, left0 + (card_w + card_gap), card_top, card_w, card_h,
              'Days 30-60', 'MySQL + auditor UX', [
                  'MySQL restore + assertion engine (Fable audit: clean split)',
                  'Auditor-view-only role (no-drill account)',
                  'Evidence-bundle: quarterly generation cron',
                  'PagerDuty AI-agent playbook integration',
                  'Goal: expand paid pilots to 25',
              ], PINK_400)
timeline_card(s, left0 + 2 * (card_w + card_gap), card_top, card_w, card_h,
              'Days 60-90', 'Enterprise motion + fundraise', [
                  'BYO-runner (deploy in customer VPC) beta',
                  'SSO (SAML + OIDC) + DPA + BAA templates',
                  'First insurance-broker partnership',
                  'Sales collateral for insurance channel',
                  'Series Seed close',
              ], TEAL_400)


# ---- save ----
out = Path(__file__).parent / 'Dokaz_Investor_Deck.pptx'
prs.save(out)
print(f'Wrote {out}')
